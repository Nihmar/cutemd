"""Convert office documents and comic archives to HTML for in-app preview."""

from __future__ import annotations

import zipfile
from pathlib import Path
from html import escape


def _wrap_html(body: str, css: str) -> str:
    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>{css}</style></head>
<body class="preview">{body}</body></html>"""


# ---------------------------------------------------------------------------
# XLSX  ->  HTML <table>
# ---------------------------------------------------------------------------

_XLSX_TABLE_CSS = """
.xlsx-toc { margin: 0.5em 0; line-height: 1.8; }
.xlsx-toc a { color: #61afef; text-decoration: none; margin-right: 1em; }
.xlsx-table { border-collapse: collapse; width: 100%; margin: 0.5em 0; font-size: 11px; }
.xlsx-table th, .xlsx-table td { border: 1px solid #777; padding: 3px 6px; text-align: left; white-space: nowrap; }
.xlsx-table th { font-weight: bold; background: #3c3c3c; }
.xlsx-table tr:nth-child(even) td { background: #2a2a2a; }
"""

def xlsx_to_html(path: Path, css: str, sheet_names: list[str] | None = None) -> str:
    """Convert an XLSX workbook to HTML table(s).

    Args:
        path: Path to the .xlsx file.
        css: Base CSS to include (appended with table-specific CSS).
        sheet_names: If not None, render only these sheets in order.
                     If None, render all sheets with a clickable TOC.
    """
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    names = sheet_names if sheet_names is not None else wb.sheetnames
    parts: list[str] = []

    # TOC when rendering all sheets
    if len(names) > 1 and sheet_names is None:
        toc_items = "".join(
            f'<a href="#sheet_{i}">{escape(n)}</a>'
            for i, n in enumerate(names)
        )
        parts.append(f'<div class="xlsx-toc">{toc_items}</div>')

    for idx, sheet_name in enumerate(names):
        if sheet_name not in wb.sheetnames:
            continue
        ws = wb[sheet_name]
        parts.append(f'<a name="sheet_{idx}"><h2>{escape(sheet_name)}</h2></a>')
        rows: list[str] = []
        for i, row in enumerate(ws.iter_rows(values_only=True, max_row=200)):
            tag = "th" if i == 0 else "td"
            cells = "".join(
                f"<{tag}>{escape(str(c)) if c is not None else ''}</{tag}>"
                for c in row
            )
            rows.append(f"<tr>{cells}</tr>")
        parts.append(f'<table class="xlsx-table">{"".join(rows)}</table>')

    wb.close()
    return _wrap_html("\n".join(parts), css + _XLSX_TABLE_CSS)


# ---------------------------------------------------------------------------
# DOCX  ->  HTML (headings, paragraphs, tables)
# ---------------------------------------------------------------------------

def _guess_mime_from_blob(data: bytes) -> str:
    if data[:8] == b'\x89PNG\r\n\x1a\n':
        return 'image/png'
    if data[:2] == b'\xff\xd8':
        return 'image/jpeg'
    if data[:6] in (b'GIF87a', b'GIF89a'):
        return 'image/gif'
    if data[:2] == b'BM':
        return 'image/bmp'
    if data[:4] == b'RIFF' and data[8:12] == b'WEBP':
        return 'image/webp'
    return 'image/png'


_DOCX_NS = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
_DOCX_A_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
_DOCX_R_NS = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'


def docx_to_html(path: Path, css: str) -> str:
    """Convert a DOCX document to HTML paragraphs with embedded images."""
    from docx import Document
    import base64

    doc = Document(str(path))

    image_map: dict[str, bytes] = {}
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_map[rel.rId] = rel.target_part.blob
                except Exception:
                    pass
    except Exception:
        pass

    parts: list[str] = []
    for el in doc.element.body:
        tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
        if tag == 'p':
            html = _docx_paragraph_to_html(el, image_map)
            if html:
                parts.append(html)
        elif tag == 'tbl':
            parts.append(_docx_table_to_html(el, image_map))

    return _wrap_html('\n'.join(parts), css)


def _docx_paragraph_to_html(p_el, image_map: dict[str, bytes]) -> str:
    import base64
    parts = []
    for r in p_el.iter(f'{_DOCX_NS}r'):
        t_el = r.find(f'{_DOCX_NS}t')
        if t_el is not None and t_el.text:
            parts.append(escape(t_el.text))
        for drawing in r.iter(f'{_DOCX_NS}drawing'):
            blip = drawing.find(f'.//{_DOCX_A_NS}blip')
            if blip is not None:
                embed = blip.get(f'{{{_DOCX_R_NS}}}embed')
                if embed and embed in image_map:
                    data = image_map[embed]
                    mime = _guess_mime_from_blob(data)
                    b64 = base64.b64encode(data).decode('ascii')
                    parts.append(f'<img src="data:{mime};base64,{b64}" style="max-width:100%" />')

    text = ''.join(parts).strip()
    if not text:
        return ''

    p_style = p_el.find(f'{_DOCX_NS}pPr')
    style_val = None
    if p_style is not None:
        p_style_el = p_style.find(f'{_DOCX_NS}pStyle')
        if p_style_el is not None:
            style_val = p_style_el.get(f'{_DOCX_NS}val')

    if style_val and style_val.startswith('Heading'):
        level = style_val.replace('Heading', '')
        if level.isdigit() and 1 <= int(level) <= 6:
            return f'<h{level}>{text}</h{level}>'

    return f'<p>{text}</p>'


def _docx_table_to_html(tbl_el, image_map: dict[str, bytes]) -> str:
    rows_html = []
    for row_el in tbl_el.iter(f'{_DOCX_NS}tr'):
        cells = []
        for cell_el in row_el.iter(f'{_DOCX_NS}tc'):
            cell_html = []
            for p in cell_el.iter(f'{_DOCX_NS}p'):
                html = _docx_paragraph_to_html(p, image_map)
                if html:
                    cell_html.append(html)
            cells.append(f'<td>{"".join(cell_html)}</td>')
        rows_html.append('<tr>' + ''.join(cells) + '</tr>')
    return '<table>' + '\n'.join(rows_html) + '</table>'


# ---------------------------------------------------------------------------
# PPTX  ->  HTML (slides with text)
# ---------------------------------------------------------------------------

def pptx_to_html(path: Path, css: str, slide_index: int | None = None) -> str:
    """Convert a PPTX presentation to HTML."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import base64

    prs = Presentation(str(path))
    parts: list[str] = []
    slides = list(prs.slides)

    if slide_index is not None:
        if 0 <= slide_index < len(slides):
            slides = [slides[slide_index]]
        else:
            return _wrap_html("<p>[Slide not found]</p>", css)

    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(f'<p>{escape(text)}</p>')
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    b64 = base64.b64encode(img.blob).decode('ascii')
                    parts.append(f'<p><img src="data:{img.content_type};base64,{b64}" style="max-width:100%" /></p>')
                except Exception:
                    pass
        parts.append('<hr>')

    return _wrap_html('\n'.join(parts), css)


# ---------------------------------------------------------------------------
# CBZ  ->  HTML (image gallery)
# ---------------------------------------------------------------------------

def cbz_to_html(path: Path, css: str) -> str:
    """Convert a CBZ (comic book zip) to HTML showing all images."""
    _IMG_EXTS = frozenset(
        {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".webp"}
    )
    images: list[str] = []
    try:
        with zipfile.ZipFile(path) as zf:
            for name in sorted(zf.namelist()):
                ext = Path(name).suffix.lower()
                if ext in _IMG_EXTS:
                    data = zf.read(name)
                    import base64
                    mime = "image/png" if ext == ".png" else "image/jpeg"
                    if ext in (".jpg", ".jpeg"):
                        mime = "image/jpeg"
                    elif ext == ".gif":
                        mime = "image/gif"
                    elif ext == ".svg":
                        mime = "image/svg+xml"
                    elif ext == ".webp":
                        mime = "image/webp"
                    b64 = base64.b64encode(data).decode("ascii")
                    images.append(f'<p><img src="data:{mime};base64,{b64}" /></p>')
    except (zipfile.BadZipFile, FileNotFoundError):
        return _wrap_html("<p>[Cannot read: file is corrupted]</p>", css)
    except Exception:
        pass

    body = "\n".join(images) if images else "<p>[No images found in CBZ file]</p>"
    style = css + """
    .cbz-gallery img {
        max-width: 100%; height: auto; display: block; margin: 1em auto;
    }"""
    return _wrap_html(f'<div class="cbz-gallery">{body}</div>', style)


# ---------------------------------------------------------------------------
# EPUB  ->  HTML (chapter text)
# ---------------------------------------------------------------------------

def _strip_html(text: str) -> str:
    """Remove HTML/XML tags, returning clean text."""
    import re
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def epub_to_html(path: Path, css: str) -> str:
    """Convert an EPUB to HTML -- text from each chapter."""
    parts: list[str] = []
    try:
        with zipfile.ZipFile(path) as zf:
            for name in sorted(zf.namelist()):
                ext = Path(name).suffix.lower()
                if ext not in (".xhtml", ".html", ".htm", ".xml"):
                    continue
                content = zf.read(name).decode("utf-8", errors="replace")
                if ext == ".xml" and "<html" not in content.lower()[:500]:
                    continue
                text = _strip_html(content)
                if not text:
                    continue
                title = Path(name).stem
                parts.append(f"<h2>{escape(title)}</h2>")
                for para in text.split("\n\n"):
                    para = para.strip()
                    if para:
                        parts.append(f"<p>{escape(para)}</p>")
    except (zipfile.BadZipFile, FileNotFoundError):
        return _wrap_html("<p>[Cannot read: file is corrupted]</p>", css)
    except Exception:
        pass

    body = "\n".join(parts[:50]) if parts else "<p>[No content found in EPUB]</p>"
    return _wrap_html(body, css)
